import os
import logging
import tempfile
from pathlib import Path
from typing import List
from io import BytesIO
from dotenv import load_dotenv
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import JSONResponse
import uvicorn

# File processing libraries
import docx2txt
from pptx import Presentation
import fitz  # PyMuPDF for PDFs

# Import text splitting utility (you may also use one from langchain or custom)
from llama_index.core import Document  # LlamaIndex's document abstraction

# Import your vector index components and FAISS vector store from LlamaIndex
from llama_index import (
    GPTSimpleVectorStoreIndex,
    ServiceContext,
    LLMPredictor,
    PromptHelper,
)
from llama_index.vector_stores import FAISSVectorStore

# For production, configure proper logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("rag_app")

load_dotenv()

GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
if not GOOGLE_API_KEY:
    raise ValueError("Google API key not set in .env file")

# --- Placeholder for Google Gemini LLM Integration ---
class GoogleGeminiLLM:
    """
    Dummy implementation of Google Gemini LLM.
    Replace this with actual API calls or library integration for gemini-2.5-pro-exp-03-25.
    """
    def __init__(self, api_key: str, model_name: str = "gemini-2.5-pro-exp-03-25"):
        self.api_key = api_key
        self.model_name = model_name
        # Initialize the Google Gemini client using the API key here.
        # For example: self.client = SomeGeminiClient(api_key=self.api_key, model=self.model_name)
    
    def predict(self, prompt: str) -> str:
        # Replace with an actual call to the Gemini service using self.client.
        # For now, just log the prompt and return a dummy response.
        print(f"GoogleGeminiLLM ({self.model_name}) received prompt: {prompt}")
        return "Processed response from Gemini."

# Then, instantiate the Gemini LLM with your API key:
gemini_llm = GoogleGeminiLLM(api_key=GOOGLE_API_KEY)

# --- Query Rewriting using Gemini LLM ---
def rewrite_query(query: str, llm: GoogleGeminiLLM) -> str:
    """
    Uses Gemini LLM to rephrase ambiguous queries.
    """
    rewrite_prompt = f"Rephrase the following query to be more specific and clear: '{query}'"
    rewritten = llm.predict(rewrite_prompt)
    logger.info(f"Rewritten query: {rewritten}")
    return rewritten

# --- Placeholder for Hybrid Retrieval ---
def perform_hybrid_retrieval(query: str, vector_index: GPTSimpleVectorIndex) -> str:
    """
    Combines BM25, dense retrieval, ColBERT scoring, and BGE reranking.
    This is a placeholder – replace with your hybrid retrieval implementation.
    """
    # You might first retrieve candidates with BM25 and FAISS dense retrieval,
    # then apply ColBERT-based late interaction scoring and BGE reranking.
    # For now, we use the vector index's query method.
    logger.info("Performing hybrid retrieval with enhanced ranking.")
    response = vector_index.query(query)
    # Optionally, process response with additional reranking logic here.
    return response.response

# --- Text Extraction Functions ---
def extract_text_from_pdf(file_bytes: bytes) -> str:
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_docx(file_bytes: bytes) -> str:
    with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    text = docx2txt.process(tmp_path)
    os.unlink(tmp_path)
    return text

def extract_text_from_pptx(file_bytes: bytes) -> str:
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    prs = Presentation(tmp_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    os.unlink(tmp_path)
    return text

def extract_text_from_txt(file_bytes: bytes) -> str:
    return file_bytes.decode("utf-8", errors="ignore")

# --- Document Processing with Chunking ---
def chunk_text(text: str, chunk_size: int = 400, overlap: int = 50) -> List[str]:
    """
    Splits the input text into chunks of about 300-500 tokens.
    This is a simple implementation; replace with a robust text splitter if needed.
    """
    # Here, we use character-based splitting as a proxy for token splitting.
    # In production, consider using a tokenizer (e.g., from HuggingFace) for precise token counts.
    tokens = text.split()
    chunks = []
    start = 0
    while start < len(tokens):
        end = min(start + chunk_size, len(tokens))
        chunk = " ".join(tokens[start:end])
        chunks.append(chunk)
        start += chunk_size - overlap  # move ahead with overlap
    return chunks

def process_file(file: UploadFile) -> List[Document]:
    file_content = file.file.read()
    text = ""
    suffix = Path(file.filename).suffix.lower()
    if suffix == ".pdf":
        text = extract_text_from_pdf(file_content)
    elif suffix in [".doc", ".docx"]:
        text = extract_text_from_docx(file_content)
    elif suffix in [".ppt", ".pptx"]:
        text = extract_text_from_pptx(file_content)
    elif suffix == ".txt":
        text = extract_text_from_txt(file_content)
    else:
        raise HTTPException(status_code=400, detail=f"Unsupported file type: {suffix}")

    # Chunk the text to improve retrieval accuracy.
    chunks = chunk_text(text, chunk_size=400, overlap=50)
    documents = [Document(chunk) for chunk in chunks if chunk.strip()]
    logger.info(f"Processed {len(documents)} document chunks from file: {file.filename}")
    return documents

# --- FastAPI App Setup ---
app = FastAPI(title="Production RAG Pipeline API")

# Global vector index; in production, consider using dependency injection or a proper state management.
vector_index: GPTSimpleVectorIndex = None

# Instantiate the Gemini LLM for query rewriting and potential LLM inference
gemini_llm = GoogleGeminiLLM()

@app.post("/training")
async def training(files: List[UploadFile] = File(...)):
    """
    Endpoint to train the RAG pipeline by uploading files.
    Processes PDFs, DOCX, TXT, and PPTX files, applies chunking,
    and builds an index using FAISS with a hybrid retriever configuration.
    """
    global vector_index
    try:
        all_documents = []
        for file in files:
            try:
                docs = process_file(file)
                all_documents.extend(docs)
            except Exception as fe:
                logger.error(f"Failed to process file {file.filename}: {fe}")
                continue

        if not all_documents:
            raise HTTPException(status_code=400, detail="No valid documents to index.")

        # Configure prompt helper and LLM predictor for indexing.
        prompt_helper = PromptHelper(max_input_size=2048, num_output=256, max_chunk_overlap=50)
        llm_predictor = LLMPredictor(llm=gemini_llm)  # Using Gemini as the LLM predictor.
        service_context = ServiceContext.from_defaults(llm_predictor=llm_predictor, prompt_helper=prompt_helper)

        # Build the FAISS-based vector index using LlamaIndex.
        # Here, GPTSimpleVectorIndex is used as a base; in production, you would further integrate
        # your hybrid retriever logic (BM25, Dense, ColBERT, BGE reranker) into the query pipeline.
        vector_index = GPTSimpleVectorIndex.from_documents(
            all_documents,
            service_context=service_context,
            vector_store_cls=FAISSVectorStore
        )
        logger.info("Indexing complete with %d documents.", len(all_documents))
        return JSONResponse(content={"status": "Training complete", "doc_count": len(all_documents)})
    except Exception as e:
        logger.exception("Error during training:")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/query")
async def query(q: str = Form(...)):
    """
    Endpoint to receive a query.
    The query is first rewritten for clarity using Gemini LLM,
    then processed by the hybrid retriever which combines BM25, Dense retrieval,
    ColBERT scoring, and BGE reranking.
    """
    global vector_index
    if vector_index is None:
        raise HTTPException(status_code=400, detail="No training data available. Please run /training first.")
    try:
        # Step 1: Rewrite query for specificity
        rewritten_query = rewrite_query(q, gemini_llm)

        # Step 2: Execute the hybrid retrieval query on the FAISS index
        answer = perform_hybrid_retrieval(rewritten_query, vector_index)

        # Return both the rewritten query and final answer for transparency
        return JSONResponse(content={"original_query": q, "rewritten_query": rewritten_query, "answer": answer})
    except Exception as e:
        logger.exception("Error during query processing:")
        raise HTTPException(status_code=500, detail=str(e))

# --- Production Uvicorn Runner ---
if __name__ == "__main__":
    # Use proper production server configuration, e.g., Gunicorn/Uvicorn workers, etc.
    uvicorn.run(app, host="0.0.0.0", port=int(os.getenv("PORT", 8000)))
